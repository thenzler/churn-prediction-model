#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Document processor for text extraction and chunking.
This module handles loading documents from various formats (PDF, DOCX, TXT),
extracting text content, and splitting into manageable chunks.
"""

import os
import re
import logging
import sys
from typing import List, Dict, Tuple, Optional, Union
import json
from pathlib import Path
import hashlib

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

try:
    import docx
    from docx.opc.exceptions import PackageNotFoundError
    DOCX_AVAILABLE = True
except ImportError:
    logger.warning("python-docx not installed. DOCX support unavailable.")
    DOCX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    logger.warning("PyPDF2 not installed. PDF support unavailable.")
    PDF_AVAILABLE = False

try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    logger.warning("python-pptx not installed. PPTX support unavailable.")
    PPTX_AVAILABLE = False

class Document:
    """Class representing a document with metadata and content."""
    
    def __init__(self, 
                 path: str, 
                 content: str = None, 
                 metadata: Dict = None, 
                 doc_type: str = None,
                 doc_id: str = None):
        """
        Initialize a document.
        
        Args:
            path: Path to the document file
            content: Text content of the document
            metadata: Dictionary of metadata
            doc_type: Type of document (pdf, docx, txt, etc.)
            doc_id: Unique identifier for the document
        """
        self.path = path
        self.filename = os.path.basename(path)
        self.content = content or ""
        self.metadata = metadata or {}
        self.doc_type = doc_type or self._detect_type()
        self.doc_id = doc_id or self._generate_id()
        self.chunks = []
        
    def _detect_type(self) -> str:
        """Detect document type from file extension."""
        _, ext = os.path.splitext(self.path)
        return ext.lower().strip(".")
    
    def _generate_id(self) -> str:
        """Generate a unique document ID."""
        # Use hash of path and last modified time for uniqueness
        try:
            mtime = os.path.getmtime(self.path)
        except (FileNotFoundError, OSError):
            mtime = 0
        content_hash = hashlib.md5(f"{self.path}:{mtime}".encode()).hexdigest()
        return content_hash
    
    def to_dict(self) -> Dict:
        """Convert document to dictionary representation."""
        return {
            "doc_id": self.doc_id,
            "filename": self.filename,
            "path": self.path,
            "doc_type": self.doc_type,
            "metadata": self.metadata,
            "content_length": len(self.content),
            "num_chunks": len(self.chunks)
        }
    
    def __repr__(self) -> str:
        return f"Document(id={self.doc_id}, filename={self.filename}, type={self.doc_type})"


class DocumentChunk:
    """Class representing a chunk of text from a document."""
    
    def __init__(self, 
                 text: str, 
                 doc_id: str, 
                 chunk_id: str, 
                 metadata: Dict = None):
        """
        Initialize a document chunk.
        
        Args:
            text: Text content of the chunk
            doc_id: ID of the parent document
            chunk_id: Unique identifier for the chunk
            metadata: Dictionary of metadata
        """
        self.text = text
        self.doc_id = doc_id
        self.chunk_id = chunk_id
        self.metadata = metadata or {}
        
    def to_dict(self) -> Dict:
        """Convert chunk to dictionary representation."""
        return {
            "chunk_id": self.chunk_id,
            "doc_id": self.doc_id,
            "text": self.text,
            "metadata": self.metadata,
            "text_length": len(self.text)
        }
    
    def __repr__(self) -> str:
        return f"DocumentChunk(id={self.chunk_id}, doc_id={self.doc_id}, len={len(self.text)})"


class DocumentProcessor:
    """Class for processing documents and extracting text."""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize document processor.
        
        Args:
            chunk_size: Maximum size of text chunks in characters
            chunk_overlap: Number of characters to overlap between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        
    def process_document(self, doc_path: str) -> Optional[Document]:
        """
        Process a document and extract text.
        
        Args:
            doc_path: Path to the document file
            
        Returns:
            Document object with extracted text, or None if processing fails
        """
        logger.info(f"Processing document: {doc_path}")
        
        if not os.path.exists(doc_path):
            logger.error(f"Document not found: {doc_path}")
            return None
        
        # Create document object
        document = Document(path=doc_path)
        
        # Extract text based on document type
        try:
            if document.doc_type == "pdf" and PDF_AVAILABLE:
                document.content = self._extract_from_pdf(doc_path)
            elif document.doc_type == "docx" and DOCX_AVAILABLE:
                document.content = self._extract_from_docx(doc_path)
            elif document.doc_type == "pptx" and PPTX_AVAILABLE:
                document.content = self._extract_from_pptx(doc_path)
            elif document.doc_type == "txt":
                document.content = self._extract_from_txt(doc_path)
            else:
                if os.path.isfile(doc_path):
                    # Try to read as plain text for unknown formats
                    document.content = self._extract_from_txt(doc_path)
                else:
                    logger.error(f"Unsupported document type: {document.doc_type}")
                    return None
            
            # Extract metadata
            document.metadata = self._extract_metadata(doc_path, document.doc_type)
            
            # Create chunks
            document.chunks = self._create_chunks(document)
            
            logger.info(f"Document processed: {document.filename} - {len(document.content)} chars, {len(document.chunks)} chunks")
            return document
            
        except Exception as e:
            logger.error(f"Error processing document: {doc_path} - {str(e)}")
            return None
    
    def process_directory(self, directory_path: str, recursive: bool = True) -> List[Document]:
        """
        Process all documents in a directory.
        
        Args:
            directory_path: Path to directory containing documents
            recursive: Whether to process subdirectories
            
        Returns:
            List of Document objects
        """
        logger.info(f"Processing directory: {directory_path}")
        
        documents = []
        
        if not os.path.exists(directory_path):
            logger.error(f"Directory not found: {directory_path}")
            return documents
        
        # Process files in directory
        for item in os.listdir(directory_path):
            item_path = os.path.join(directory_path, item)
            
            if os.path.isfile(item_path):
                # Process file extensions we support
                _, ext = os.path.splitext(item_path)
                if ext.lower() in ['.pdf', '.docx', '.txt', '.pptx']:
                    doc = self.process_document(item_path)
                    if doc is not None:
                        documents.append(doc)
            
            elif os.path.isdir(item_path) and recursive:
                # Process subdirectory if recursive is True
                subdirectory_docs = self.process_directory(item_path, recursive)
                documents.extend(subdirectory_docs)
        
        logger.info(f"Processed {len(documents)} documents in directory: {directory_path}")
        return documents
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF document."""
        if not PDF_AVAILABLE:
            raise ImportError("PyPDF2 is not installed. Cannot process PDF files.")
        
        text = ""
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n\n"
        
        return text
    
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX document."""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx is not installed. Cannot process DOCX files.")
        
        try:
            doc = docx.Document(file_path)
            text = "\n\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text
        except PackageNotFoundError:
            raise ValueError(f"Invalid DOCX file: {file_path}")
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX document."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not installed. Cannot process PPTX files.")
        
        presentation = pptx.Presentation(file_path)
        text = ""
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n\n"
        
        return text
    
    def _extract_from_txt(self, file_path: str) -> str:
        """Extract text from TXT document."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try different encodings if UTF-8 fails
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read()
            except:
                logger.error(f"Failed to read text file with multiple encodings: {file_path}")
                return ""
    
    def _extract_metadata(self, file_path: str, doc_type: str) -> Dict:
        """Extract metadata from document."""
        metadata = {
            "filename": os.path.basename(file_path),
            "file_size": os.path.getsize(file_path),
            "created_time": os.path.getctime(file_path),
            "modified_time": os.path.getmtime(file_path)
        }
        
        # Extract additional metadata based on document type
        if doc_type == "pdf" and PDF_AVAILABLE:
            try:
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    if pdf_reader.metadata:
                        for key, value in pdf_reader.metadata.items():
                            # Remove the leading slash from PDF metadata keys
                            clean_key = key.strip("/") if isinstance(key, str) else key
                            metadata[clean_key] = value
                    metadata["page_count"] = len(pdf_reader.pages)
            except:
                logger.warning(f"Failed to extract PDF metadata from {file_path}")
        
        elif doc_type == "docx" and DOCX_AVAILABLE:
            try:
                doc = docx.Document(file_path)
                core_properties = doc.core_properties
                metadata["title"] = core_properties.title
                metadata["author"] = core_properties.author
                metadata["created"] = core_properties.created
                metadata["modified"] = core_properties.modified
                metadata["last_modified_by"] = core_properties.last_modified_by
            except:
                logger.warning(f"Failed to extract DOCX metadata from {file_path}")
        
        return metadata
    
    def _create_chunks(self, document: Document) -> List[DocumentChunk]:
        """Split document content into chunks."""
        text = document.content
        chunks = []
        
        # Simple splitting by characters with overlap
        start = 0
        chunk_counter = 0
        
        while start < len(text):
            # If we're near the end, just include the remainder
            if start + self.chunk_size >= len(text):
                end = len(text)
            else:
                # Try to find a suitable breakpoint (newline or period)
                ideal_end = start + self.chunk_size
                
                # Look for a period or newline in the nearby area
                look_ahead = min(len(text), ideal_end + 100)
                
                # Try to find a paragraph break first (double newline)
                paragraph_break = text.find('\n\n', ideal_end, look_ahead)
                if paragraph_break != -1:
                    end = paragraph_break + 2
                else:
                    # Try to find a single newline
                    newline = text.find('\n', ideal_end, look_ahead)
                    if newline != -1:
                        end = newline + 1
                    else:
                        # Try to find a period followed by space
                        period = text.find('. ', ideal_end, look_ahead)
                        if period != -1:
                            end = period + 2
                        else:
                            # If no natural break point, use the chunk size
                            end = ideal_end
            
            # Extract chunk text
            chunk_text = text[start:end]
            
            # Only create a chunk if it has content
            if chunk_text.strip():
                chunk_id = f"{document.doc_id}_{chunk_counter}"
                
                # Create metadata for the chunk
                chunk_metadata = {
                    "filename": document.filename,
                    "doc_path": document.path,
                    "chunk_index": chunk_counter,
                    "char_start": start,
                    "char_end": end
                }
                
                # Create a new chunk object
                chunk = DocumentChunk(
                    text=chunk_text,
                    doc_id=document.doc_id,
                    chunk_id=chunk_id,
                    metadata=chunk_metadata
                )
                
                chunks.append(chunk)
                chunk_counter += 1
            
            # Move to the next chunk, accounting for overlap
            start = end - self.chunk_overlap
            # Make sure we're making forward progress
            if start <= 0 or start >= len(text):
                break
        
        return chunks
    
    def save_processed_documents(self, documents: List[Document], output_dir: str):
        """
        Save processed documents and chunks to disk.
        
        Args:
            documents: List of Document objects
            output_dir: Directory to save processed documents
        """
        # Create output directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)
        
        # Create subdirectories
        docs_dir = os.path.join(output_dir, "documents")
        chunks_dir = os.path.join(output_dir, "chunks")
        os.makedirs(docs_dir, exist_ok=True)
        os.makedirs(chunks_dir, exist_ok=True)
        
        # Save document index
        doc_index = [doc.to_dict() for doc in documents]
        with open(os.path.join(output_dir, "document_index.json"), 'w') as f:
            json.dump(doc_index, f, indent=2)
        
        # Save each document's metadata and content
        for doc in documents:
            doc_path = os.path.join(docs_dir, f"{doc.doc_id}.json")
            with open(doc_path, 'w') as f:
                json.dump({
                    "metadata": doc.to_dict(),
                    "content": doc.content
                }, f, indent=2)
            
            # Save chunks for this document
            for chunk in doc.chunks:
                chunk_path = os.path.join(chunks_dir, f"{chunk.chunk_id}.json")
                with open(chunk_path, 'w') as f:
                    json.dump(chunk.to_dict(), f, indent=2)
        
        logger.info(f"Saved {len(documents)} documents with {sum(len(doc.chunks) for doc in documents)} chunks to {output_dir}")

def main():
    """Main function to process documents."""
    import argparse
    
    parser = argparse.ArgumentParser(description='Process documents for question answering')
    parser.add_argument('--input', required=True, help='Path to document or directory')
    parser.add_argument('--output', default='processed_documents', help='Output directory for processed documents')
    parser.add_argument('--chunk-size', type=int, default=1000, help='Size of document chunks in characters')
    parser.add_argument('--chunk-overlap', type=int, default=200, help='Overlap between chunks in characters')
    parser.add_argument('--recursive', action='store_true', help='Process directories recursively')
    
    args = parser.parse_args()
    
    # Initialize processor
    processor = DocumentProcessor(
        chunk_size=args.chunk_size,
        chunk_overlap=args.chunk_overlap
    )
    
    # Process input
    if os.path.isfile(args.input):
        documents = [processor.process_document(args.input)]
        documents = [doc for doc in documents if doc is not None]
    elif os.path.isdir(args.input):
        documents = processor.process_directory(args.input, args.recursive)
    else:
        logger.error(f"Input path not found: {args.input}")
        return 1
    
    # Save processed documents
    if documents:
        processor.save_processed_documents(documents, args.output)
        logger.info(f"Successfully processed {len(documents)} documents.")
    else:
        logger.warning("No documents were successfully processed.")
    
    return 0

if __name__ == "__main__":
    sys.exit(main())
